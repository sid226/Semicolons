
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from Models import slides_res
from Utils import utils


class Slides:
    def __init__(self, path):
        self.path = path

    def create_ppt(self, title):
        filename = self.deck_title(title)
        self.title_content(f"PPT/{filename}", "Task Status")
        res = slides_res.SlidesRes(
            f"/{utils.BASE_ENDPOINT}/v1/download-ppt/{filename}")
        return res, 201

    def deck_title(self, ppt_title):

        # create presentation
        prs = Presentation(self.path)
        # add title slide
        deck_slide_layout = prs.slide_layouts[0]
        deck_slide = prs.slides.add_slide(deck_slide_layout)
        title = deck_slide.shapes.title
        subtitle = deck_slide.placeholders[1]
        title.text = ppt_title
        subtitle.text = "Feb 17, 2023\nSprint 1"
        filename = "SPRINT1-02172023.pptx"
        prs.save(f"PPT/{filename}")

        return filename

    def title_content(self, filepath, ppt_title):
        prs = Presentation(filepath)

        content_slide_layout = prs.slide_layouts[3]
        content_slide = prs.slides.add_slide(content_slide_layout)
        title = content_slide.shapes.title
        title.text = ppt_title

        # Add a bullet point paragraph
        bullet_paragraph = content_slide.shapes.placeholders[1].text_frame.add_paragraph(
        )
        bullet_paragraph.text = "Completed"
        bullet_paragraph.font.bold = True
        bullet_paragraph.level = 0

        completed_task = [
            "Alston has done most of the development work related to add cart functionality and identified/fixed a defect related to login functionality.",
            "Bob has completed the task of adding tests for the signup page and identified/fixed a defect related to OTP during signup.",
            "Jaz has created a database for order details and the backend team has started using it."
        ]
        for item in completed_task:
            self.display_points(content_slide.shapes, item)

        # Add a bullet point paragraph
        bullet_paragraph = content_slide.shapes.placeholders[1].text_frame.add_paragraph(
        )
        bullet_paragraph.text = "In-Progress"
        bullet_paragraph.font.bold = True
        bullet_paragraph.level = 0

        inprogress = [
            "Alston is raising a PR for the work done on Jira 1412 and looking into defects related to Jira 3452.",
            "Bob is working on adding test cases for the login page and is currently blocked due to issues in fetching data from the mock database.",
            "Jaz is currently resolving the defect related to frequent DB connection failures on Jira 3232."
        ]
        for item in inprogress:
            self.display_points(content_slide.shapes, item)

        prs.save(filepath)

    def display_points(self,shapes, item):
        bullet_paragraph = shapes.placeholders[1].text_frame.add_paragraph()
        bullet_paragraph.text = item
        bullet_paragraph.font.bold = False
        bullet_paragraph.level = 1
